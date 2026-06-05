import os
import asyncio
from datetime import datetime
# pyrefly: ignore [missing-import]
from fastapi import FastAPI, HTTPException, WebSocket, WebSocketDisconnect, Form
# pyrefly: ignore [missing-import]
from fastapi.middleware.cors import CORSMiddleware
# pyrefly: ignore [missing-import]
from fastapi.responses import StreamingResponse
# pyrefly: ignore [missing-import]
from fastapi.staticfiles import StaticFiles
# pyrefly: ignore [missing-import]
from pydantic import BaseModel  
from openpyxl import Workbook
from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
from openpyxl.utils import get_column_letter
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors
# pyrefly: ignore [missing-import]
from pptx import Presentation
# pyrefly: ignore [missing-import]
from pptx.util import Inches, Pt
# pyrefly: ignore [missing-import]
from pptx.dml.color import RGBColor
from io import BytesIO
from typing import List, Dict, Any, Optional,Set
from app.agents.sql_agent import agent_executor, db_manager, vector_store

app = FastAPI(
    title="Enterprise AI SQL Agent API",
    description="Backend API supporting conversational query execution, AST validation, and PII masking.",
    version="1.0.0"
)

# Enable CORS for local client development
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

class QueryRequest(BaseModel):
    query: str
    role: str = "general"

class QueryResponse(BaseModel):
    user_query: str
    user_role: str
    relevant_tables: List[str]
    generated_sql: Optional[str]
    query_results: Optional[List[Dict[str, Any]]]
    execution_error: Optional[str]
    retry_count: int
    narrative_response: str

class GlossaryTermRequest(BaseModel):
    term: str
    definition: str
    sql_hint: str

class ExcelExportRequest(BaseModel): 
    results: List[Dict[str, Any]]

class PDFExportRequest(BaseModel): 
    query: str
    narrative: str
    results: List[Dict[str, Any]]

class PPTXExportRequest(BaseModel): 
    query: str
    narrative: str
    sql: str 


@app.post("/api/v1/query", response_model=QueryResponse)
async def run_query(request: QueryRequest):
    inputs = {
        "user_query": request.query,
        "user_role": request.role,
        "relevant_tables": [],
        "table_schemas": "",
        "glossary_terms": [],
        "generated_sql": None,
        "query_results": None,
        "execution_error": None,
        "retry_count": 0,
        "narrative_response": None
    }
    try:
        final_state = await agent_executor.ainvoke(inputs)
        return QueryResponse(
            user_query=final_state["user_query"],
            user_role=final_state["user_role"],
            relevant_tables=final_state["relevant_tables"],
            generated_sql=final_state["generated_sql"],
            query_results=final_state["query_results"],
            execution_error=final_state["execution_error"],
            retry_count=final_state["retry_count"],
            narrative_response=final_state["narrative_response"] or ""
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Agent execution error: {str(e)}")

@app.get("/api/v1/tables")
async def get_tables():
    try:
        tables = await db_manager.get_table_names()
        return {"tables": tables}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to fetch table list: {str(e)}")

@app.get("/api/v1/tables/{table_name}/schema")
async def get_table_schema(table_name: str):
    try:
        tables = await db_manager.get_table_names()
        if table_name not in tables:
            raise HTTPException(status_code=404, detail=f"Table '{table_name}' not found.")
        
        columns = await db_manager.get_table_schema(table_name)
        fkeys = await db_manager.get_foreign_keys(table_name)
        return {
            "table": table_name,
            "columns": columns,
            "foreign_keys": fkeys
        }
    except HTTPException as he:
        raise he
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to fetch table schema: {str(e)}")

@app.get("/api/v1/glossary")
async def search_glossary(q: str = ""):
    try:
        results = vector_store.search_glossary(q, limit=10)
        return {"results": results}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Glossary lookup failed: {str(e)}")

@app.post("/api/v1/glossary")
async def add_glossary_term(request: GlossaryTermRequest):
    try:
        vector_store.upsert_glossary_term(request.term, request.definition, request.sql_hint)
        return {"status": "success", "message": f"Glossary term '{request.term}' registered."}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to register glossary term: {str(e)}")

# --- WEBSOCKET Room Session Manager ---
class ConnectionManager: 
    def __init__(self): 
        self.active_connections: Dict[str, Set[WebSocket]] = {}
    
    async def connect(self, websocket: WebSocket, room_id: str): 
        await websocket.accept()

        if room_id not in self.active_connections: 
            self.active_connections[room_id] = set()
        self.active_connections[room_id].add(websocket)
    
    def disconnect(self, websocket: WebSocket, room_id: str): 
        if room_id in self.active_connections: 
            self.active_connections[room_id].remove(websocket)
            if not self.active_connections[room_id]: 
                del self.active_connections[room_id]
    
    async def broadcast(self, message: dict, room_id: str, exclude_websocket: WebSocket = None): 
        if room_id in self.active_connections: 
            for connection in self.active_connections[room_id]: 
                if connection != exclude_websocket: 
                    try: 
                        await connection.send_json(message)
                    except Exception: 
                        pass

manager = ConnectionManager()

@app.websocket("/ws/warroom/{room_id}")
async def websocket_endpoint(websocket: WebSocket, room_id: str): 
    await manager.connect(websocket, room_id)
    try: 
        while True: 
            data = await websocket.receive_json()
            await manager.broadcast(data, room_id, exclude_websocket=websocket)
    
    except WebSocketDisconnect: 
        manager.disconnect(websocket, room_id)


# --- Document Exporter Routes ---

@app.post("/api/v1/export/excel")
async def export_excel(req: ExcelExportRequest): 
    if not req.results: 
        raise HTTPException(status_code=400, detail="No data available to export.")
    
    wb = Workbook()
    ws = wb.active
    ws.title = "Query Results"

    headers = list(req.results[0].keys())
    ws.append(headers)

    # Header styling
    header_font = Font(name="Segoe UI", size=11, bold=True, color="FFFFFF")
    header_fill = PatternFill(start_color="1E293B", end_color="1E293B", fill_type="solid")
    center_align = Alignment(horizontal="center", vertical="center")
    border_side = Side(border_style="thin", color="E2E8F0")
    cell_border = Border(left=border_side, right=border_side, top=border_side, bottom=border_side)

    for col_num in range(1, len(headers) + 1): 
        cell = ws.cell(row=1, column=col_num)
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = center_align
        cell.border = cell_border
        
    for row_data in req.results:
        row_values = [row_data[h] for h in headers]
        ws.append(row_values)
    
    #Append a SUM formula summary row for numerical columns
    row_count = len(req.results)
    if row_count > 0:
        summary_row = ["Total Summary"] + [""] * (len(headers) - 1)
        ws.append(summary_row)
        ws.cell(row=row_count + 2, column=1).font = Font(bold=True)
        
        for idx, header in enumerate(headers):
            first_val = req.results[0][header]
            if isinstance(first_val, (int, float)) and idx > 0:
                col_letter = get_column_letter(idx + 1)
                formula_cell = ws.cell(row=row_count + 2, column=idx + 1)
                formula_cell.value = f"=SUM({col_letter}2:{col_letter}{row_count + 1})"
                formula_cell.font = Font(bold=True)
                formula_cell.border = cell_border
                
    for col in ws.columns:
        max_len = max(len(str(cell.value or '')) for cell in col)
        col_letter = get_column_letter(col[0].column)
        ws.column_dimensions[col_letter].width = max(max_len + 3, 12)
    
    file_stream = BytesIO()
    wb.save(file_stream)
    file_stream.seek(0)
    
    return StreamingResponse(
        file_stream,
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers={"Content-Disposition": "attachment; filename=query_results.xlsx"}
    )

@app.post("/api/v1/export/pptx")
async def export_pptx(req: PPTXExportRequest):
    prs = Presentation()
    
    slate_dark = RGBColor(15, 23, 42)
    blue_primary = RGBColor(37, 99, 235)
    text_dark = RGBColor(51, 65, 85)
    
    # Slide 1: Cover Slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = "SQL Intelligence Slide Deck"
    p.font.bold = True
    p.font.size = Pt(36)
    p.font.color.rgb = blue_primary
    
    query_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(2))
    qtf = query_box.text_frame
    qp = qtf.add_paragraph()
    qp.text = f"User Request:\n\"{req.query}\""
    qp.font.italic = True
    qp.font.size = Pt(18)
    qp.font.color.rgb = text_dark
    
    # Slide 2: SQL Statement
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = "Generated SQL Query"
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = blue_primary
    
    sql_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(4.5))
    sq_tf = sql_box.text_frame
    sq_tf.word_wrap = True
    sq_p = sq_tf.add_paragraph()
    sq_p.text = req.sql
    sq_p.font.name = "Courier New"
    sq_p.font.size = Pt(12)
    sq_p.font.color.rgb = slate_dark
    
    # Slide 3: Executive Narrative
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = "Executive Briefing Summary"
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = blue_primary
    
    narr_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(4.5))
    n_tf = narr_box.text_frame
    n_tf.word_wrap = True
    
    for paragraph_text in req.narrative.split("\n"):
        if paragraph_text.strip():
            np = n_tf.add_paragraph()
            np.text = paragraph_text.strip()
            np.font.size = Pt(14)
            np.font.color.rgb = text_dark
            np.space_after = Pt(8)
            
    file_stream = BytesIO()
    prs.save(file_stream)
    file_stream.seek(0)
    
    return StreamingResponse(
        file_stream,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": "attachment; filename=briefing_deck.pptx"}
    )

@app.post("/api/v1/export/pdf")
async def export_pdf(req: PDFExportRequest):
    file_stream = BytesIO()
    doc = SimpleDocTemplate(file_stream, pagesize=letter, rightMargin=40, leftMargin=40, topMargin=40, bottomMargin=40)
    story = []
    
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        'DocTitle',
        parent=styles['Heading1'],
        fontName='Helvetica-Bold',
        fontSize=18,
        textColor=colors.HexColor('#0F172A'),
        spaceAfter=12
    )
    section_style = ParagraphStyle(
        'DocSection',
        parent=styles['Heading2'],
        fontName='Helvetica-Bold',
        fontSize=12,
        textColor=colors.HexColor('#2563EB'),
        spaceBefore=12,
        spaceAfter=6
    )
    body_style = ParagraphStyle(
        'DocBody',
        parent=styles['BodyText'],
        fontName='Helvetica',
        fontSize=10,
        leading=14,
        textColor=colors.HexColor('#334155'),
        spaceAfter=10
    )
    
    story.append(Paragraph("AI SQL Agent Intelligence Report", title_style))
    story.append(Paragraph(f"<b>Query Focus:</b> <i>\"{req.query}\"</i>", body_style))
    story.append(Spacer(1, 10))
    
    story.append(Paragraph("Executive Narrative Brief", section_style))
    story.append(Paragraph(req.narrative.replace("\n", "<br/>"), body_style))
    story.append(Spacer(1, 10))
    
    if req.results:
        story.append(Paragraph("Database Grid Extract (First 25 Rows)", section_style))
        headers = list(req.results[0].keys())
        data_table_content = [headers]
        
        for row in req.results[:25]:
            data_table_content.append([str(row[h] if row[h] is not None else 'NULL') for h in headers])
            
        table = Table(data_table_content, colWidths=[(500 / len(headers))] * len(headers))
        table.setStyle(TableStyle([
            ('BACKGROUND', (0,0), (-1,0), colors.HexColor('#1E293B')),
            ('TEXTCOLOR', (0,0), (-1,0), colors.whitesmoke),
            ('ALIGN', (0,0), (-1,-1), 'LEFT'),
            ('FONTNAME', (0,0), (-1,0), 'Helvetica-Bold'),
            ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor('#E2E8F0')),
            ('ROWBACKGROUNDS', (0,1), (-1,-1), [colors.white, colors.HexColor('#F8FAFC')]),
            ('BOTTOMPADDING', (0,0), (-1,-1), 4),
            ('TOPPADDING', (0,0), (-1,-1), 4),
            ('FONTSIZE', (0,0), (-1,-1), 8),
        ]))
        story.append(table)
        
    doc.build(story)
    file_stream.seek(0)
    
    return StreamingResponse(
        file_stream,
        media_type="application/pdf",
        headers={"Content-Disposition": "attachment; filename=analytics_report.pdf"}
    )

# --- Alerts Registry & Smoke Detector Background Scheduler ---

ALERTS_REGISTRY = [
    {
        "id": "alert_1",
        "name": "High Invoice Volumes",
        "query": "SELECT COUNT(*) FROM Invoice;",
        "condition": "> 400",
        "interval_seconds": 15,
        "last_checked": None,
        "status": "Active"
    },
    {
        "id": "alert_2",
        "name": "High Support Ticket Count",
        "query": "SELECT COUNT(*) FROM Customer WHERE Country = 'Brazil';",
        "condition": "> 4",
        "interval_seconds": 30,
        "last_checked": None,
        "status": "Active"
    }
]

ALERTS_LOGS = []

class CreateAlertRequest(BaseModel):
    name: str
    query: str
    condition: str
    interval_seconds: int

async def alert_scheduler_loop():
    # Wait for database manager to be ready
    await asyncio.sleep(3)
    while True:
        now = datetime.now()
        for alert in ALERTS_REGISTRY:
            if alert["status"] == "Active":
                try:
                    # Run background SQL check query
                    results = await db_manager.execute_query(alert["query"])
                    alert["last_checked"] = now.strftime("%Y-%m-%d %H:%M:%S")
                    
                    if results:
                        first_row = results[0]
                        first_val = list(first_row.values())[0]
                        
                        # Evaluate condition: e.g. "> 100", "< 5", "= 10"
                        cond = alert["condition"].strip()
                        triggered = False
                        try:
                            val_num = float(first_val)
                            if cond.startswith(">"):
                                threshold = float(cond.replace(">", "").strip())
                                triggered = val_num > threshold
                            elif cond.startswith("<"):
                                threshold = float(cond.replace("<", "").strip())
                                triggered = val_num < threshold
                            elif cond.startswith("="):
                                threshold = float(cond.replace("=", "").strip())
                                triggered = val_num == threshold
                        except (ValueError, TypeError):
                            # Fallback to string comparison if not numeric
                            if cond.startswith("="):
                                threshold = cond.replace("=", "").strip()
                                triggered = str(first_val) == threshold
                        
                        if triggered:
                            alert["status"] = "Triggered"
                            ALERTS_LOGS.append({
                                "id": f"log_{len(ALERTS_LOGS) + 1}",
                                "alert_id": alert["id"],
                                "name": alert["name"],
                                "timestamp": now.strftime("%Y-%m-%d %H:%M:%S"),
                                "value": str(first_val),
                                "message": f"Check value {first_val} triggered condition '{cond}'."
                            })
                except Exception as e:
                    alert["status"] = "Error"
                    print(f"Error checking alert '{alert['name']}': {e}")
                    
        await asyncio.sleep(5)

@app.on_event("startup")
async def startup_event():
    asyncio.create_task(alert_scheduler_loop())


# --- Alerts Endpoints ---

@app.get("/api/v1/alerts")
async def get_alerts():
    return {"alerts": ALERTS_REGISTRY}

@app.post("/api/v1/alerts")
async def create_alert(req: CreateAlertRequest):
    new_alert = {
        "id": f"alert_{len(ALERTS_REGISTRY) + 1}",
        "name": req.name,
        "query": req.query,
        "condition": req.condition,
        "interval_seconds": req.interval_seconds,
        "last_checked": None,
        "status": "Active"
    }
    ALERTS_REGISTRY.append(new_alert)
    return {"status": "success", "alert": new_alert}

@app.get("/api/v1/alerts/logs")
async def get_alerts_logs():
    return {"logs": ALERTS_LOGS}

@app.post("/api/v1/alerts/{alert_id}/reset")
async def reset_alert(alert_id: str):
    for alert in ALERTS_REGISTRY:
        if alert["id"] == alert_id:
            alert["status"] = "Active"
            return {"status": "success", "message": f"Alert '{alert['name']}' reset to Active."}
    raise HTTPException(status_code=404, detail="Alert rule not found.")


# --- Slack Webhook slash command receiver ---

@app.post("/api/v1/webhooks/slack")
async def slack_webhook(text: str = Form(...)):
    inputs = {
        "user_query": text,
        "user_role": "general",
        "relevant_tables": [],
        "table_schemas": "",
        "glossary_terms": [],
        "generated_sql": None,
        "query_results": None,
        "execution_error": None,
        "retry_count": 0,
        "narrative_response": None
    }
    try:
        final_state = await agent_executor.ainvoke(inputs)
        sql = final_state.get("generated_sql") or "-- No SQL statements executed."
        narrative = final_state.get("narrative_response") or "No insights compiled."
        
        # Slack Block Kit payload structure
        blocks = [
            {
                "type": "header",
                "text": {
                    "type": "plain_text",
                    "text": "📊 AI SQL Agent Intelligence Briefing",
                    "emoji": True
                }
            },
            {
                "type": "section",
                "text": {
                    "type": "mrkdwn",
                    "text": f"*Query Request:* \"{text}\""
                }
            },
            {
                "type": "section",
                "text": {
                    "type": "mrkdwn",
                    "text": f"*Insight Summary:*\n{narrative}"
                }
            },
            {
                "type": "context",
                "elements": [
                    {
                        "type": "mrkdwn",
                        "text": f"*Compiled SQL:* `{sql}`"
                    }
                ]
            }
        ]
        return {"response_type": "in_channel", "blocks": blocks}
    except Exception as e:
        return {
            "response_type": "ephemeral",
            "text": f"⚠️ Agent execution failed: {str(e)}"
        }


# Mount static frontend files to root path
frontend_dir = os.path.normpath(
    os.path.join(os.path.dirname(os.path.abspath(__file__)), "../../../frontend/dist")
)

if os.path.exists(frontend_dir):
    app.mount("/", StaticFiles(directory=frontend_dir, html=True), name="static")

