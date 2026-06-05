from io import BytesIO
# pyrefly: ignore [missing-import]
from fastapi import APIRouter, HTTPException
# pyrefly: ignore [missing-import]
from fastapi.responses import StreamingResponse
from openpyxl import Workbook   
from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
from openpyxl.utils import get_column_letter
# pyrefly: ignore [missing-import]
from pptx import Presentation
# pyrefly: ignore [missing-import]
from pptx.util import Inches, Pt
# pyrefly: ignore [missing-import]
from pptx.dml.color import RGBColor
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors

from app.api.schemas import ExcelExportRequest, PDFExportRequest, PPTXExportRequest

router = APIRouter(prefix="/api/v1/export", tags=["Exports"])

@router.post("/excel")
async def export_excel(req: ExcelExportRequest): 
    if not req.results: 
        raise HTTPException(status_code=400, detail="No data available to export.")
    
    wb = Workbook()
    ws = wb.active
    ws.title = "Query Results"

    headers = list(req.results[0].keys())
    ws.append(headers)

    # Header styling
    header_font = Font(name="Segoe UI", size=11, bold=True, color="FFFFFF")
    header_fill = PatternFill(start_color="1E293B", end_color="1E293B", fill_type="solid")
    center_align = Alignment(horizontal="center", vertical="center")
    border_side = Side(border_style="thin", color="E2E8F0")
    cell_border = Border(left=border_side, right=border_side, top=border_side, bottom=border_side)

    for col_num in range(1, len(headers) + 1): 
        cell = ws.cell(row=1, column=col_num)
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = center_align
        cell.border = cell_border
        
    for row_data in req.results:
        row_values = [row_data[h] for h in headers]
        ws.append(row_values)
    
    # Append a SUM formula summary row for numerical columns
    row_count = len(req.results)
    if row_count > 0:
        summary_row = ["Total Summary"] + [""] * (len(headers) - 1)
        ws.append(summary_row)
        ws.cell(row=row_count + 2, column=1).font = Font(bold=True)

        for idx, header in enumerate(headers):
            first_val = req.results[0][header]
            if isinstance(first_val, (int, float)) and idx > 0:
                col_letter = get_column_letter(idx + 1)
                formula_cell = ws.cell(row=row_count + 2, column=idx + 1)
                formula_cell.value = f"=SUM({col_letter}2:{col_letter}{row_count + 1})"
                formula_cell.font = Font(bold=True)
                formula_cell.border = cell_border
                
    for col in ws.columns:
        max_len = max(len(str(cell.value or '')) for cell in col)
        col_letter = get_column_letter(col[0].column)
        ws.column_dimensions[col_letter].width = max(max_len + 3, 12)
    
    file_stream = BytesIO()
    wb.save(file_stream)
    file_stream.seek(0)
    
    return StreamingResponse(
        file_stream,
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers={"Content-Disposition": "attachment; filename=query_results.xlsx"}
    )

@router.post("/pptx")
async def export_pptx(req: PPTXExportRequest):
    prs = Presentation()
    
    slate_dark = RGBColor(15, 23, 42)
    blue_primary = RGBColor(37, 99, 235)
    text_dark = RGBColor(51, 65, 85)
    
    # Slide 1: Cover Slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = "SQL Intelligence Slide Deck"
    p.font.bold = True
    p.font.size = Pt(36)
    p.font.color.rgb = blue_primary
    
    query_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(2))
    qtf = query_box.text_frame
    qp = qtf.add_paragraph()
    qp.text = f"User Request:\n\"{req.query}\""
    qp.font.italic = True
    qp.font.size = Pt(18)
    qp.font.color.rgb = text_dark
    
    # Slide 2: SQL Statement
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = "Generated SQL Query"
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = blue_primary
    
    sql_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(4.5))
    sq_tf = sql_box.text_frame
    sq_tf.word_wrap = True
    sq_p = sq_tf.add_paragraph()
    sq_p.text = req.sql
    sq_p.font.name = "Courier New"
    sq_p.font.size = Pt(12)
    sq_p.font.color.rgb = slate_dark
    
    # Slide 3: Executive Narrative
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = "Executive Briefing Summary"
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = blue_primary
    
    narr_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(4.5))
    n_tf = narr_box.text_frame
    n_tf.word_wrap = True
    
    for paragraph_text in req.narrative.split("\n"):
        if paragraph_text.strip():
            np = n_tf.add_paragraph()
            np.text = paragraph_text.strip()
            np.font.size = Pt(14)
            np.font.color.rgb = text_dark
            np.space_after = Pt(8)
            
    file_stream = BytesIO()
    prs.save(file_stream)
    file_stream.seek(0)
    
    return StreamingResponse(
        file_stream,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": "attachment; filename=briefing_deck.pptx"}
    )

@router.post("/pdf")
async def export_pdf(req: PDFExportRequest):
    file_stream = BytesIO()
    doc = SimpleDocTemplate(file_stream, pagesize=letter, rightMargin=40, leftMargin=40, topMargin=40, bottomMargin=40)
    story = []
    
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        'DocTitle',
        parent=styles['Heading1'],
        fontName='Helvetica-Bold',
        fontSize=18,
        textColor=colors.HexColor('#0F172A'),
        spaceAfter=12
    )
    section_style = ParagraphStyle(
        'DocSection',
        parent=styles['Heading2'],
        fontName='Helvetica-Bold',
        fontSize=12,
        textColor=colors.HexColor('#2563EB'),
        spaceBefore=12,
        spaceAfter=6
    )
    body_style = ParagraphStyle(
        'DocBody',
        parent=styles['BodyText'],
        fontName='Helvetica',
        fontSize=10,
        leading=14,
        textColor=colors.HexColor('#334155'),
        spaceAfter=10
    )
    
    story.append(Paragraph("AI SQL Agent Intelligence Report", title_style))
    story.append(Paragraph(f"<b>Query Focus:</b> <i>\"{req.query}\"</i>", body_style))
    story.append(Spacer(1, 10))
    
    story.append(Paragraph("Executive Narrative Brief", section_style))
    story.append(Paragraph(req.narrative.replace("\n", "<br/>"), body_style))
    story.append(Spacer(1, 10))
    
    if req.results:
        story.append(Paragraph("Database Grid Extract (First 25 Rows)", section_style))
        headers = list(req.results[0].keys())
        data_table_content = [headers]
        
        for row in req.results[:25]:
            data_table_content.append([str(row[h] if row[h] is not None else 'NULL') for h in headers])
            
        table = Table(data_table_content, colWidths=[(500 / len(headers))] * len(headers))
        table.setStyle(TableStyle([
            ('BACKGROUND', (0,0), (-1,0), colors.HexColor('#1E293B')),
            ('TEXTCOLOR', (0,0), (-1,0), colors.whitesmoke),
            ('ALIGN', (0,0), (-1,-1), 'LEFT'),
            ('FONTNAME', (0,0), (-1,0), 'Helvetica-Bold'),
            ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor('#E2E8F0')),
            ('ROWBACKGROUNDS', (0,1), (-1,-1), [colors.white, colors.HexColor('#F8FAFC')]),
            ('BOTTOMPADDING', (0,0), (-1,-1), 4),
            ('TOPPADDING', (0,0), (-1,-1), 4),
            ('FONTSIZE', (0,0), (-1,-1), 8),
        ]))
        story.append(table)
        
    doc.build(story)
    file_stream.seek(0)
    
    return StreamingResponse(
        file_stream,
        media_type="application/pdf",
        headers={"Content-Disposition": "attachment; filename=analytics_report.pdf"}
    )
